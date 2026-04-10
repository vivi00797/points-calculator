#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import xml.etree.ElementTree as ET

# Parse the XMind content.xml
tree = ET.parse('content.xml')
root = tree.getroot()

# Define namespace
ns = {'xmind': 'urn:xmind:xmap:xmlns:content:2.0'}

# Extract main topics
main_topics = []
root_topic = root.find('.//xmind:topic', ns)
if root_topic is not None:
    title_elem = root_topic.find('xmind:title', ns)
    main_title = title_elem.text if title_elem is not None else "通信原理"
    
    # Get all main topics (children of root)
    children = root_topic.find('.//xmind:children/xmind:topics[@type="attached"]', ns)
    if children is not None:
        for topic in children.findall('xmind:topic', ns):
            title_elem = topic.find('xmind:title', ns)
            if title_elem is not None:
                topic_title = title_elem.text
                
                # Get subtopics
                subtopics = []
                sub_children = topic.find('.//xmind:children/xmind:topics[@type="attached"]', ns)
                if sub_children is not None:
                    for subtopic in sub_children.findall('xmind:topic', ns):
                        sub_title_elem = subtopic.find('xmind:title', ns)
                        if sub_title_elem is not None:
                            subtopics.append(sub_title_elem.text)
                
                main_topics.append({
                    'title': topic_title,
                    'subtopics': subtopics[:8]  # Limit to 8 subtopics per slide
                })

# Create PowerPoint presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define a simple color scheme
TITLE_COLOR = RGBColor(31, 78, 120)
SUBTITLE_COLOR = RGBColor(68, 114, 196)
TEXT_COLOR = RGBColor(64, 64, 64)

def add_title_slide(prs, title):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

def add_toc_slide(prs, topics):
    """Add table of contents slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "目录"
    
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Add topics
    left = Inches(1.5)
    top = Inches(1.8)
    width = Inches(7)
    height = Inches(5)
    
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, topic in enumerate(topics, 1):
        p = content_frame.add_paragraph() if i > 1 else content_frame.paragraphs[0]
        p.text = f"{i}. {topic['title']}"
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(12)

def add_content_slide(prs, topic_data):
    """Add content slide for each main topic"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = topic_data['title']
    
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Add subtopics if any
    if topic_data['subtopics']:
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(7)
        height = Inches(5)
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, subtopic in enumerate(topic_data['subtopics']):
            p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
            p.text = f"• {subtopic}"
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(10)
            p.level = 0

# Create presentation
print(f"Creating presentation with {len(main_topics)} main topics...")

# Title slide
add_title_slide(prs, main_title)

# Table of contents
add_toc_slide(prs, main_topics)

# Content slides
for topic in main_topics:
    add_content_slide(prs, topic)

# Save presentation
output_file = '通信原理.pptx'
prs.save(output_file)
print(f"Presentation saved as: {output_file}")
print(f"Total slides: {len(prs.slides)}")
