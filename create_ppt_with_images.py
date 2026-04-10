#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import xml.etree.ElementTree as ET
import os
from pathlib import Path

# Parse the XMind content.xml
tree = ET.parse('content.xml')
root = tree.getroot()

# Define namespace
ns = {'xmind': 'urn:xmind:xmap:xmlns:content:2.0'}

# Extract main topics
main_topics = []
root_topic = root.find('.//xmind:topic', ns)
if root_topic is not None:
    title_elem = root_topic.find('xmind:title', ns)
    main_title = title_elem.text if title_elem is not None else "通信原理"
    
    # Get all main topics (children of root)
    children = root_topic.find('.//xmind:children/xmind:topics[@type="attached"]', ns)
    if children is not None:
        for topic in children.findall('xmind:topic', ns):
            title_elem = topic.find('xmind:title', ns)
            if title_elem is not None:
                topic_title = title_elem.text
                
                # Get subtopics
                subtopics = []
                sub_children = topic.find('.//xmind:children/xmind:topics[@type="attached"]', ns)
                if sub_children is not None:
                    for subtopic in sub_children.findall('xmind:topic', ns):
                        sub_title_elem = subtopic.find('xmind:title', ns)
                        if sub_title_elem is not None:
                            subtopics.append(sub_title_elem.text)
                
                main_topics.append({
                    'title': topic_title,
                    'subtopics': subtopics[:8]  # Limit to 8 subtopics per slide
                })

# Image mapping for topics
image_dir = Path.home() / '.gemini/antigravity/brain/f2834e40-e9dd-4175-8967-4ed2bf8a252b'
topic_images = {
    '绪论': 'communication_intro',
    '确知信号': 'deterministic_signals',
    '随机过程': 'random_process',
    '信道': 'communication_channel',
    '模拟调制系统': 'analog_modulation',
    '数字基带传输系统': 'digital_baseband',
}

# Create PowerPoint presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define a simple color scheme
TITLE_COLOR = RGBColor(31, 78, 120)
SUBTITLE_COLOR = RGBColor(68, 114, 196)
TEXT_COLOR = RGBColor(64, 64, 64)
ACCENT_COLOR = RGBColor(91, 155, 213)

def add_title_slide(prs, title):
    """Add title slide with decorative elements"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add decorative shape
    left = Inches(0)
    top = Inches(2.5)
    width = Inches(10)
    height = Inches(2.5)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.color.rgb = ACCENT_COLOR
    
    # Add title
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_toc_slide(prs, topics):
    """Add table of contents slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add decorative header bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.color.rgb = ACCENT_COLOR
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "目录"
    
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add topics in two columns
    left_col = Inches(1)
    right_col = Inches(5.5)
    top = Inches(2)
    width = Inches(4)
    height = Inches(5)
    
    left_box = slide.shapes.add_textbox(left_col, top, width, height)
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    right_box = slide.shapes.add_textbox(right_col, top, width, height)
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    mid_point = (len(topics) + 1) // 2
    
    for i, topic in enumerate(topics[:mid_point], 1):
        p = left_frame.add_paragraph() if i > 1 else left_frame.paragraphs[0]
        p.text = f"{i}. {topic['title']}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(10)
    
    for i, topic in enumerate(topics[mid_point:], mid_point + 1):
        p = right_frame.add_paragraph() if i > mid_point + 1 else right_frame.paragraphs[0]
        p.text = f"{i}. {topic['title']}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(10)

def find_image_file(image_name):
    """Find image file in the artifacts directory"""
    for ext in ['.png', '.jpg', '.jpeg']:
        pattern = f"{image_name}_*{ext}"
        files = list(image_dir.glob(pattern))
        if files:
            return str(files[0])
    return None

def add_content_slide(prs, topic_data, topic_index):
    """Add content slide for each main topic with image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add decorative header bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.color.rgb = ACCENT_COLOR
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = topic_data['title']
    
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Try to add image if available
    image_added = False
    if topic_data['title'] in topic_images:
        image_file = find_image_file(topic_images[topic_data['title']])
        if image_file and os.path.exists(image_file):
            left = Inches(0.5)
            top = Inches(1.5)
            height = Inches(3)
            slide.shapes.add_picture(image_file, left, top, height=height)
            image_added = True
            content_left = Inches(4.5)
            content_width = Inches(5)
        else:
            # Add decorative shape if no image
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(3)
            height = Inches(3)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(220, 230, 241)
            shape.line.color.rgb = ACCENT_COLOR
            shape.line.width = Pt(2)
            content_left = Inches(4)
            content_width = Inches(5.5)
    else:
        # No image, use decorative shape
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(3)
        height = Inches(3)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(220, 230, 241)
        shape.line.color.rgb = ACCENT_COLOR
        shape.line.width = Pt(2)
        
        # Add topic number in the shape
        text_frame = shape.text_frame
        text_frame.text = str(topic_index)
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(120)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        text_frame.vertical_anchor = 1  # Middle
        
        content_left = Inches(4)
        content_width = Inches(5.5)
    
    # Add subtopics
    if topic_data['subtopics']:
        top = Inches(1.8)
        height = Inches(5)
        
        content_box = slide.shapes.add_textbox(content_left, top, content_width, height)
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, subtopic in enumerate(topic_data['subtopics']):
            p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
            p.text = f"• {subtopic}"
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(8)
            p.level = 0

# Create presentation
print(f"Creating presentation with {len(main_topics)} main topics...")

# Title slide
add_title_slide(prs, main_title)

# Table of contents
add_toc_slide(prs, main_topics)

# Content slides
for idx, topic in enumerate(main_topics, 1):
    add_content_slide(prs, topic, idx)

# Save presentation
output_file = '通信原理.pptx'
prs.save(output_file)
print(f"Presentation saved as: {output_file}")
print(f"Total slides: {len(prs.slides)}")
print(f"Images added for topics: {list(topic_images.keys())}")
